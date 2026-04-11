"""
Generate Mid-Semester PPT for AD Attack-Path Discovery Mapper
Uses python-pptx to create a professional, well-aligned presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# =============================================================================
# CONFIGURATION
# =============================================================================
OUTPUT_FILE = os.path.join(os.path.dirname(__file__), "AD_Path_Discovery_Mapper_MidSem.pptx")

# Color Palette (dark professional cyber theme)
BG_DARK        = RGBColor(0x0D, 0x11, 0x17)   # Slide background
BG_CARD        = RGBColor(0x14, 0x1B, 0x2D)   # Card / content box
ACCENT_BLUE    = RGBColor(0x00, 0xB4, 0xD8)   # Headings / accents
ACCENT_CYAN    = RGBColor(0x48, 0xCA, 0xE4)   # Sub-accents
ACCENT_GREEN   = RGBColor(0x06, 0xD6, 0xA0)   # Highlights
ACCENT_ORANGE  = RGBColor(0xFF, 0x92, 0x3E)   # Warnings / attention
ACCENT_RED     = RGBColor(0xEF, 0x47, 0x6F)   # Critical
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY     = RGBColor(0xBB, 0xBB, 0xCC)
MID_GRAY       = RGBColor(0x88, 0x88, 0x99)
DARK_GRAY      = RGBColor(0x22, 0x2A, 0x3A)

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(tf, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment

def add_paragraph(tf, text, font_size=16, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(4), font_name="Calibri", level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    return p

def add_slide_number(slide, num, total):
    tb = add_text_box(slide, Inches(11.8), Inches(7.0), Inches(1.4), Inches(0.4))
    set_text(tb.text_frame, f"{num} / {total}", font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

def add_footer_line(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.9), Inches(12.333), Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    shape.shadow.inherit = False

def add_section_header(slide, title):
    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.35), Inches(0.25), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()
    bar.shadow.inherit = False
    # Title text
    tb = add_text_box(slide, Inches(1.0), Inches(0.25), Inches(11.0), Inches(0.7))
    set_text(tb.text_frame, title, font_size=28, bold=True, color=ACCENT_BLUE, font_name="Calibri")

TOTAL_SLIDES = 12

# =============================================================================
# SLIDE 1 — TITLE SLIDE
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# Accent top bar
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_BLUE)

# Center card
card = add_rect(slide, Inches(2.0), Inches(0.8), Inches(9.333), Inches(5.8), BG_CARD, ACCENT_BLUE, Pt(1.5))

# "Mid Semester Presentation" label
tb = add_text_box(slide, Inches(2.5), Inches(1.1), Inches(8.333), Inches(0.5))
set_text(tb.text_frame, "MID SEMESTER PRESENTATION", font_size=14, bold=True, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER, font_name="Calibri")

# Title
tb = add_text_box(slide, Inches(2.5), Inches(1.7), Inches(8.333), Inches(1.2))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Active Directory Attack-Path", font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "Discovery Mapper", font_size=36, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER, space_before=Pt(2))

# Divider
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.15), Inches(2.333), Pt(2))
div.fill.solid()
div.fill.fore_color.rgb = ACCENT_BLUE
div.line.fill.background()
div.shadow.inherit = False

# Submitted to
tb = add_text_box(slide, Inches(2.5), Inches(3.4), Inches(8.333), Inches(0.7))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Submitted to", font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "School of Cyber Security & Digital Forensics", font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, space_before=Pt(2))
add_paragraph(tf, "National Forensic Sciences University", font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(2))

# Student info
tb = add_text_box(slide, Inches(2.5), Inches(4.6), Inches(8.333), Inches(0.9))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Student Name  |  Enrollment No.  |  M.Tech Cyber Security & Incident Response", font_size=14, bold=False, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

# Guide info
tb = add_text_box(slide, Inches(2.5), Inches(5.3), Inches(8.333), Inches(0.9))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Under the Guidance of", font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "Guide Name  |  Designation", font_size=14, bold=False, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER, space_before=Pt(4))

# Bottom bar
add_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), ACCENT_BLUE)

add_slide_number(slide, 1, TOTAL_SLIDES)

# =============================================================================
# SLIDE 2 — TABLE OF CONTENTS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Table of Contents")
add_footer_line(slide)

toc_items = [
    ("01", "Abstract"),
    ("02", "Literature Review"),
    ("03", "Motivation"),
    ("04", "Problem Statement"),
    ("05", "Project Objectives"),
    ("06", "Methodology & Approach"),
    ("07", "Architecture & System Design"),
    ("08", "Implementation — Lab Setup & Tools"),
    ("09", "Implementation — Features & Demo"),
    ("10", "Next Course of Action"),
    ("11", "Guide Consultation Schedule"),
    ("12", "Thank You & Questions"),
]

# Two-column layout
col_left_x  = Inches(1.0)
col_right_x = Inches(7.0)
start_y      = Inches(1.3)
row_h        = Inches(0.44)

for i, (num, title) in enumerate(toc_items):
    col_x = col_left_x if i < 6 else col_right_x
    y_pos = start_y + (i % 6) * row_h

    # Number box
    nb = add_rect(slide, col_x, y_pos, Inches(0.55), Inches(0.35), ACCENT_BLUE)
    nb.text_frame.paragraphs[0].text = num
    nb.text_frame.paragraphs[0].font.size = Pt(13)
    nb.text_frame.paragraphs[0].font.bold = True
    nb.text_frame.paragraphs[0].font.color.rgb = WHITE
    nb.text_frame.paragraphs[0].font.name = "Calibri"
    nb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    nb.text_frame.word_wrap = False

    # Title
    ttb = add_text_box(slide, col_x + Inches(0.7), y_pos, Inches(4.8), Inches(0.35))
    set_text(ttb.text_frame, title, font_size=15, color=WHITE, alignment=PP_ALIGN.LEFT)

add_slide_number(slide, 2, TOTAL_SLIDES)

# =============================================================================
# SLIDE 3 — ABSTRACT
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Abstract")
add_footer_line(slide)

abstract_text = (
    "Active Directory (AD) is the backbone of identity and access management in over 90% of Fortune 500 "
    "enterprises. However, misconfigurations in AD — such as overly permissive Access Control Lists (ACLs), "
    "excessive group nesting, Kerberoastable service accounts, and unrestricted delegation — create hidden "
    "privilege escalation paths that adversaries routinely exploit during post-compromise operations."
)
abstract_text2 = (
    "The \"AD Attack-Path Discovery Mapper\" is a preventive security tool that performs authenticated LDAP "
    "enumeration against a Windows Domain Controller, constructs a directed permission graph of all principals "
    "and their relationships using NetworkX, and applies BFS/DFS algorithms to discover attack paths from "
    "low-privileged users to high-value targets such as Domain Admins."
)
abstract_text3 = (
    "The tool provides a modern web dashboard built with Next.js featuring interactive Cytoscape.js graph "
    "visualization, severity-classified findings, MITRE ATT&CK technique tagging, and copy-ready PowerShell "
    "remediation scripts. Unlike existing tools like BloodHound that focus on post-exploitation analysis, this "
    "project operates strictly from a defender's perspective — emphasizing proactive misconfiguration detection "
    "and prioritized remediation guidance without relying on SIEMs, EDR telemetry, or Windows event logs."
)

card = add_rect(slide, Inches(0.7), Inches(1.2), Inches(11.933), Inches(5.5), BG_CARD, DARK_GRAY, Pt(1))

tb = add_text_box(slide, Inches(1.1), Inches(1.4), Inches(11.2), Inches(5.0))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, abstract_text, font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.JUSTIFY)
add_paragraph(tf, "", font_size=8, color=LIGHT_GRAY, alignment=PP_ALIGN.JUSTIFY)
add_paragraph(tf, abstract_text2, font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.JUSTIFY, space_before=Pt(8))
add_paragraph(tf, "", font_size=8, color=LIGHT_GRAY, alignment=PP_ALIGN.JUSTIFY)
add_paragraph(tf, abstract_text3, font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.JUSTIFY, space_before=Pt(8))

# Key tags at bottom
tags = ["LDAP Enumeration", "Graph Analysis", "MITRE ATT&CK", "PowerShell Remediation", "Preventive Security"]
tag_x = Inches(1.1)
for tag in tags:
    w = Inches(len(tag) * 0.12 + 0.3)
    t = add_rect(slide, tag_x, Inches(5.9), w, Inches(0.32), DARK_GRAY, ACCENT_BLUE, Pt(1))
    t.text_frame.paragraphs[0].text = tag
    t.text_frame.paragraphs[0].font.size = Pt(9)
    t.text_frame.paragraphs[0].font.color.rgb = ACCENT_CYAN
    t.text_frame.paragraphs[0].font.name = "Calibri"
    t.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tag_x += w + Inches(0.15)

add_slide_number(slide, 3, TOTAL_SLIDES)

# =============================================================================
# SLIDE 4 — LITERATURE REVIEW
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Literature Review")
add_footer_line(slide)

papers = [
    {
        "id": "1",
        "authors": "Dunagan et al. (Microsoft Research, 2009)",
        "title": "Heat-Ray: Combating Identity Snowball Attacks Using Machine Learning",
        "findings": "Introduced graph-based modeling of credential relationships in enterprise networks; demonstrated ML-based detection of lateral movement chains.",
        "limitations": "Focused on credential hopping; did not address ACL-based privilege escalation paths or provide automated remediation guidance.",
    },
    {
        "id": "2",
        "authors": "Robbins et al. (SpecterOps, 2017)",
        "title": "BloodHound — Six Degrees of Domain Admin",
        "findings": "Pioneered the use of graph theory (Neo4j) to map AD attack paths from any user to Domain Admin; widely adopted by red teams globally.",
        "limitations": "Requires SharpHound data collector running on a domain-joined host; designed as an offensive tool with limited defensive remediation workflow; no MITRE ATT&CK tagging.",
    },
    {
        "id": "3",
        "authors": "Metcalf, S. (ADSecurity.org, 2018)",
        "title": "Active Directory Security Risk Assessment & Attack Path Analysis",
        "findings": "Catalogued 15+ common AD misconfigurations (unconstrained delegation, SPN abuse, AdminCount) and their exploitation techniques.",
        "limitations": "Provided manual audit checklists; lacked automated graph-based path discovery or a unified dashboard for continuous monitoring.",
    },
    {
        "id": "4",
        "authors": "Alsaheel et al. (IEEE S&P, 2021)",
        "title": "ATLAS: A Sequence-Based Framework for Enterprise Threat Analysis",
        "findings": "Used audit logs and NLP to reconstruct attack sequences in enterprise networks; achieved high accuracy in identifying multi-stage attacks.",
        "limitations": "Relies heavily on Windows Event Log telemetry and SIEM data; cannot operate in environments without centralized logging infrastructure.",
    },
    {
        "id": "5",
        "authors": "Guo et al. (ACSAC, 2022)",
        "title": "ProGrapher: Proactive Detection of Lateral Movement Using Graph Embeddings",
        "findings": "Applied graph neural networks to authentication graphs for early detection of lateral movement; reduced false positives compared to rule-based approaches.",
        "limitations": "Requires substantial training data from compromised environments; focused on detection rather than prevention or remediation; high computational overhead.",
    },
]

start_y = Inches(1.2)
for i, paper in enumerate(papers):
    y = start_y + i * Inches(1.14)
    # Numbering circle
    circ = add_rect(slide, Inches(0.6), y + Inches(0.05), Inches(0.35), Inches(0.35), ACCENT_BLUE)
    circ.text_frame.paragraphs[0].text = paper["id"]
    circ.text_frame.paragraphs[0].font.size = Pt(12)
    circ.text_frame.paragraphs[0].font.bold = True
    circ.text_frame.paragraphs[0].font.color.rgb = WHITE
    circ.text_frame.paragraphs[0].font.name = "Calibri"
    circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Content card
    card = add_rect(slide, Inches(1.1), y, Inches(11.6), Inches(1.05), BG_CARD, DARK_GRAY, Pt(0.75))
    tb = add_text_box(slide, Inches(1.25), y + Inches(0.03), Inches(11.3), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, f"{paper['authors']} — \"{paper['title']}\"", font_size=12, bold=True, color=ACCENT_CYAN, alignment=PP_ALIGN.LEFT)
    add_paragraph(tf, f"Findings: {paper['findings']}", font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.JUSTIFY, space_before=Pt(2), space_after=Pt(1))
    add_paragraph(tf, f"Limitations: {paper['limitations']}", font_size=10, color=ACCENT_ORANGE, alignment=PP_ALIGN.JUSTIFY, space_before=Pt(1), space_after=Pt(1))

add_slide_number(slide, 4, TOTAL_SLIDES)

# =============================================================================
# SLIDE 5 — MOTIVATION
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Motivation")
add_footer_line(slide)

motivations = [
    ("🔓", "80% of Breaches Involve AD", 
     "Microsoft's 2023 Digital Defense Report states that identity-based attacks targeting Active Directory account for over 80% of enterprise breaches, with misconfigured ACLs being the #1 root cause."),
    ("📰", "SolarWinds & Kaseya Attacks", 
     "High-profile supply chain attacks (SolarWinds 2020, Kaseya 2021) exploited AD trust relationships and group policy misconfigurations to achieve domain dominance — gaps that proactive path analysis could have detected."),
    ("🛡", "Defender's Tooling Gap", 
     "Existing tools like BloodHound and PingCastle are powerful but offense-oriented. Organizations lack an integrated, defender-friendly dashboard that combines path discovery with MITRE mapping and actionable remediation scripts."),
    ("🎓", "Academic Research Gap", 
     "Literature focuses on detection (post-breach) rather than prevention (pre-breach). No open-source tool combines LDAP enumeration, graph-based path analysis, misconfiguration detection, and automated remediation in a single defensive framework."),
]

for i, (icon, title, desc) in enumerate(motivations):
    y = Inches(1.3) + i * Inches(1.4)
    card = add_rect(slide, Inches(0.7), y, Inches(11.933), Inches(1.25), BG_CARD, DARK_GRAY, Pt(0.75))
    
    # Icon box
    ib = add_rect(slide, Inches(0.9), y + Inches(0.15), Inches(0.6), Inches(0.55), ACCENT_BLUE)
    ib.text_frame.paragraphs[0].text = icon
    ib.text_frame.paragraphs[0].font.size = Pt(20)
    ib.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Text
    tb = add_text_box(slide, Inches(1.7), y + Inches(0.08), Inches(10.7), Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, title, font_size=16, bold=True, color=ACCENT_CYAN, alignment=PP_ALIGN.LEFT)
    add_paragraph(tf, desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.JUSTIFY, space_before=Pt(4))

add_slide_number(slide, 5, TOTAL_SLIDES)

# =============================================================================
# SLIDE 6 — PROBLEM STATEMENT
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Problem Statement")
add_footer_line(slide)

# Main question box
qcard = add_rect(slide, Inches(1.5), Inches(1.3), Inches(10.333), Inches(1.5), BG_CARD, ACCENT_BLUE, Pt(1.5))
tb = add_text_box(slide, Inches(1.8), Inches(1.4), Inches(9.8), Inches(1.3))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Core Research Question", font_size=12, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "\"How can defenders proactively discover and remediate privilege escalation paths in Active Directory before adversaries exploit them — without relying on SIEM, EDR, or event log telemetry?\"", font_size=17, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER, space_before=Pt(8))

problems = [
    ("Hidden Attack Paths", "ACL misconfigurations create invisible chains (User→Group→DA) that native AD tools do not visualize or audit."),
    ("Reactive Security Posture", "Current tools detect attacks after exploitation. No open-source tool provides proactive, continuous path analysis with remediation scripts."),
    ("Manual Audit Complexity", "AD environments with 500+ objects make manual permission review infeasible. Graph-based automation is needed for scalable analysis."),
]

for i, (title, desc) in enumerate(problems):
    y = Inches(3.2) + i * Inches(1.2)
    # Red indicator
    ind = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), y + Inches(0.1), Inches(0.15), Inches(0.8))
    ind.fill.solid(); ind.fill.fore_color.rgb = ACCENT_RED
    ind.line.fill.background(); ind.shadow.inherit = False

    card = add_rect(slide, Inches(1.8), y, Inches(10.033), Inches(1.0), BG_CARD, DARK_GRAY, Pt(0.75))
    tb = add_text_box(slide, Inches(2.0), y + Inches(0.05), Inches(9.6), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, title, font_size=15, bold=True, color=ACCENT_RED, alignment=PP_ALIGN.LEFT)
    add_paragraph(tf, desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.JUSTIFY, space_before=Pt(3))

add_slide_number(slide, 6, TOTAL_SLIDES)

# =============================================================================
# SLIDE 7 — PROJECT OBJECTIVES
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Project Objectives")
add_footer_line(slide)

objectives = [
    ("01", "Automated AD Reconnaissance via LDAP",
     "Perform authenticated LDAP/LDAPS enumeration to extract all Users, Groups, Computers, OUs, ACLs, SPNs, and trust relationships from a Windows Domain Controller — without agents, scripts on the DC, or event log access."),
    ("02", "Graph-Based Attack Path Discovery",
     "Construct a directed permission graph using NetworkX and apply BFS/DFS traversal algorithms to discover all privilege escalation paths from any low-privileged principal to high-value targets (Domain Admins, Enterprise Admins)."),
    ("03", "Misconfiguration Detection & Risk Scoring",
     "Identify 15+ categories of AD security misconfigurations (Kerberoastable admins, unconstrained delegation, ACL abuse, password policy violations) and compute a domain-wide risk score (0–100)."),
    ("04", "MITRE ATT&CK Mapping & Remediation",
     "Map every discovered finding to corresponding MITRE ATT&CK techniques (T1078, T1098, T1558, etc.) and auto-generate copy-ready PowerShell remediation scripts for each finding."),
    ("05", "Interactive Web Dashboard",
     "Deliver a modern, responsive web dashboard with Cytoscape.js graph visualization, object explorer, attack path viewer, findings panel, and exportable PDF/CSV reports for stakeholder communication."),
]

for i, (num, title, desc) in enumerate(objectives):
    y = Inches(1.2) + i * Inches(1.18)
    # Number
    nb = add_rect(slide, Inches(0.6), y + Inches(0.08), Inches(0.5), Inches(0.4), ACCENT_BLUE)
    nb.text_frame.paragraphs[0].text = num
    nb.text_frame.paragraphs[0].font.size = Pt(14)
    nb.text_frame.paragraphs[0].font.bold = True
    nb.text_frame.paragraphs[0].font.color.rgb = WHITE
    nb.text_frame.paragraphs[0].font.name = "Calibri"
    nb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    card = add_rect(slide, Inches(1.25), y, Inches(11.45), Inches(1.05), BG_CARD, DARK_GRAY, Pt(0.75))
    tb = add_text_box(slide, Inches(1.45), y + Inches(0.05), Inches(11.0), Inches(0.95))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, title, font_size=14, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.LEFT)
    add_paragraph(tf, desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.JUSTIFY, space_before=Pt(3))

add_slide_number(slide, 7, TOTAL_SLIDES)

# =============================================================================
# SLIDE 8 — METHODOLOGY / APPROACH
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Methodology & Approach")
add_footer_line(slide)

# Pipeline stages as horizontal flow
stages = [
    ("Phase 1", "LDAP\nEnumeration", ACCENT_BLUE),
    ("Phase 2", "Data\nNormalization", ACCENT_CYAN),
    ("Phase 3", "Graph\nConstruction", ACCENT_GREEN),
    ("Phase 4", "Path Analysis\n& Misconfigs", ACCENT_ORANGE),
    ("Phase 5", "MITRE Mapping\n& Remediation", ACCENT_RED),
    ("Phase 6", "Dashboard\nVisualization", RGBColor(0xA8, 0x5C, 0xFF)),
]

stage_w = Inches(1.7)
gap = Inches(0.28)
total_w = len(stages) * stage_w + (len(stages) - 1) * gap
start_x = (SLIDE_W - total_w) / 2

for i, (phase, label, color) in enumerate(stages):
    x = start_x + i * (stage_w + gap)
    y = Inches(1.4)

    box = add_rect(slide, x, y, stage_w, Inches(1.1), BG_CARD, color, Pt(1.5))
    box.text_frame.word_wrap = True
    set_text(box.text_frame, phase, font_size=10, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_paragraph(box.text_frame, label, font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, space_before=Pt(4))

    # Arrow (except last)
    if i < len(stages) - 1:
        arrow_x = x + stage_w
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, y + Inches(0.35), gap, Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = color
        arrow.line.fill.background()
        arrow.shadow.inherit = False

# Description cards below
method_details = [
    ("Data Collection", "Authenticated LDAPS connection (port 636) to Domain Controller using ldap3 library. Enumerates Users, Groups, Computers, OUs, ACLs, GPOs, Trusts, and SPNs via targeted LDAP queries."),
    ("Graph Modeling", "Constructs a directed graph using NetworkX. Nodes represent AD principals (Users, Groups, Computers). Edges represent relationships (MemberOf, AdminTo, HasACL, TrustedBy) with weighted severity scores."),
    ("Analysis Engine", "BFS/DFS traversal discovers all paths from any principal to Domain Admins. Misconfig detector checks 15+ patterns. Risk scoring aggregates finding severity into a 0–100 domain risk score."),
    ("Output & Remediation", "Findings mapped to MITRE ATT&CK techniques. PowerShell remediation auto-generated per finding. Results rendered on interactive Cytoscape.js graph with filtering, search, and export capabilities."),
]

for i, (title, desc) in enumerate(method_details):
    col = i % 2
    row = i // 2
    x = Inches(0.7) + col * Inches(6.3)
    y = Inches(3.0) + row * Inches(1.85)
    card = add_rect(slide, x, y, Inches(6.0), Inches(1.65), BG_CARD, DARK_GRAY, Pt(0.75))
    tb = add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.6), Inches(1.45))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, title, font_size=14, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.LEFT)
    add_paragraph(tf, desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.JUSTIFY, space_before=Pt(4))

add_slide_number(slide, 8, TOTAL_SLIDES)

# =============================================================================
# SLIDE 9 — ARCHITECTURE (GRAPHICAL)
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "System Architecture & Tools")
add_footer_line(slide)

# Architecture diagram using shapes
# Three-tier: AD Environment → Backend (Python) → Frontend (Next.js)

# --- AD Environment ---
ad_box = add_rect(slide, Inches(0.5), Inches(1.8), Inches(3.0), Inches(4.5), BG_CARD, ACCENT_RED, Pt(1.5))
tb = add_text_box(slide, Inches(0.5), Inches(1.4), Inches(3.0), Inches(0.4))
set_text(tb.text_frame, "AD Environment (MARVEL.local)", font_size=12, bold=True, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)

ad_components = ["Domain Controller", "LDAPS (Port 636)", "Users & Groups", "Computers & OUs", "ACLs & GPOs", "Trust Relationships"]
for i, comp in enumerate(ad_components):
    cb = add_rect(slide, Inches(0.8), Inches(2.0) + i * Inches(0.65), Inches(2.4), Inches(0.5), DARK_GRAY, ACCENT_RED, Pt(0.75))
    cb.text_frame.paragraphs[0].text = comp
    cb.text_frame.paragraphs[0].font.size = Pt(11)
    cb.text_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY
    cb.text_frame.paragraphs[0].font.name = "Calibri"
    cb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Arrow AD → Backend
arr1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.6), Inches(3.5), Inches(0.8), Inches(0.5))
arr1.fill.solid(); arr1.fill.fore_color.rgb = ACCENT_BLUE
arr1.line.fill.background(); arr1.shadow.inherit = False
tb = add_text_box(slide, Inches(3.55), Inches(3.1), Inches(0.9), Inches(0.35))
set_text(tb.text_frame, "LDAPS", font_size=9, bold=True, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# --- Backend ---
be_box = add_rect(slide, Inches(4.5), Inches(1.8), Inches(4.2), Inches(4.5), BG_CARD, ACCENT_BLUE, Pt(1.5))
tb = add_text_box(slide, Inches(4.5), Inches(1.4), Inches(4.2), Inches(0.4))
set_text(tb.text_frame, "Backend (Python / FastAPI)", font_size=12, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

be_modules = [
    ("ldap_enum.py", "LDAP Enumeration"),
    ("normalizer.py", "Data Normalization"),
    ("graph_builder.py", "NetworkX Graph"),
    ("attack_paths.py", "BFS/DFS Analysis"),
    ("misconfig_detector.py", "Misconfig Detection"),
    ("mitre_mapping.py", "MITRE ATT&CK"),
    ("remediation.py", "PowerShell Gen"),
]
for i, (fname, desc) in enumerate(be_modules):
    cb = add_rect(slide, Inches(4.75), Inches(1.95) + i * Inches(0.57), Inches(3.7), Inches(0.45), DARK_GRAY, ACCENT_BLUE, Pt(0.75))
    cb.text_frame.word_wrap = True
    cb.text_frame.paragraphs[0].text = f"{fname}  —  {desc}"
    cb.text_frame.paragraphs[0].font.size = Pt(10)
    cb.text_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY
    cb.text_frame.paragraphs[0].font.name = "Calibri"
    cb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Arrow Backend → Frontend
arr2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.8), Inches(3.5), Inches(0.8), Inches(0.5))
arr2.fill.solid(); arr2.fill.fore_color.rgb = ACCENT_GREEN
arr2.line.fill.background(); arr2.shadow.inherit = False
tb = add_text_box(slide, Inches(8.75), Inches(3.1), Inches(0.9), Inches(0.35))
set_text(tb.text_frame, "REST API", font_size=9, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

# --- Frontend ---
fe_box = add_rect(slide, Inches(9.7), Inches(1.8), Inches(3.1), Inches(4.5), BG_CARD, ACCENT_GREEN, Pt(1.5))
tb = add_text_box(slide, Inches(9.7), Inches(1.4), Inches(3.1), Inches(0.4))
set_text(tb.text_frame, "Frontend (Next.js / React)", font_size=12, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

fe_components = ["Dashboard Overview", "Object Explorer", "Attack Path Viewer", "Graph Visualization", "Findings & Risk Panel", "Remediation Tracker"]
for i, comp in enumerate(fe_components):
    cb = add_rect(slide, Inches(9.95), Inches(2.0) + i * Inches(0.65), Inches(2.6), Inches(0.5), DARK_GRAY, ACCENT_GREEN, Pt(0.75))
    cb.text_frame.paragraphs[0].text = comp
    cb.text_frame.paragraphs[0].font.size = Pt(11)
    cb.text_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY
    cb.text_frame.paragraphs[0].font.name = "Calibri"
    cb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Tools legend at bottom
tools_text = "Python 3.11  •  FastAPI  •  ldap3  •  NetworkX  •  Next.js 14  •  React  •  TypeScript  •  Cytoscape.js  •  TailwindCSS  •  Pydantic"
tb = add_text_box(slide, Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.35))
set_text(tb.text_frame, tools_text, font_size=11, bold=False, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 9, TOTAL_SLIDES)

# =============================================================================
# SLIDE 10 — IMPLEMENTATION (Lab Setup + Features)
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Implementation — Lab Setup & Features")
add_footer_line(slide)

# Lab Setup (left column)
lab_card = add_rect(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5), BG_CARD, DARK_GRAY, Pt(0.75))
tb = add_text_box(slide, Inches(0.7), Inches(1.3), Inches(5.4), Inches(0.4))
set_text(tb.text_frame, "Lab Environment Setup", font_size=16, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.LEFT)

lab_items = [
    ("Host OS:", "Windows 10/11 (Attacker Machine)"),
    ("Hypervisor:", "VMware / VirtualBox"),
    ("Domain Controller:", "Windows Server 2019 (192.168.186.152)"),
    ("Domain:", "MARVEL.local"),
    ("Protocol:", "LDAPS over Port 636 (TLS)"),
    ("Users:", "7 domain accounts (Tony Stark, Natasha Romanoff, etc.)"),
    ("Groups:", "~45 groups incl. Domain Admins, AVENGERS-IT"),
    ("Computers:", "3 machines (DC, WORKSTATION01, WORKSTATION02)"),
    ("ACLs:", "WriteDACL, GenericAll, GenericWrite, WriteOwner"),
    ("Python:", "3.11 with FastAPI, ldap3, NetworkX"),
    ("Frontend:", "Next.js 14 + React + TypeScript + Cytoscape.js"),
]

tb = add_text_box(slide, Inches(0.7), Inches(1.8), Inches(5.4), Inches(4.8))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "", font_size=1, color=BG_CARD)
for label, value in lab_items:
    p = add_paragraph(tf, "", font_size=12, color=WHITE, alignment=PP_ALIGN.LEFT, space_before=Pt(3), space_after=Pt(1))
    run1 = p.add_run()
    run1.text = f"  {label}  "
    run1.font.size = Pt(12)
    run1.font.bold = True
    run1.font.color.rgb = ACCENT_CYAN
    run1.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = value
    run2.font.size = Pt(12)
    run2.font.color.rgb = LIGHT_GRAY
    run2.font.name = "Calibri"

# Features (right column)
feat_card = add_rect(slide, Inches(6.7), Inches(1.2), Inches(6.1), Inches(5.5), BG_CARD, DARK_GRAY, Pt(0.75))
tb = add_text_box(slide, Inches(6.9), Inches(1.3), Inches(5.7), Inches(0.4))
set_text(tb.text_frame, "Implemented Features", font_size=16, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.LEFT)

features = [
    "✅  LDAPS authentication with NTLM (zero credential storage)",
    "✅  Full LDAP enumeration — Users, Groups, Computers, OUs, ACLs",
    "✅  Directed permission graph via NetworkX",
    "✅  BFS/DFS attack path discovery engine",
    "✅  15+ misconfiguration detection rules",
    "✅  Domain-wide risk score (0–100)",
    "✅  MITRE ATT&CK technique tagging (T1078, T1098, T1558…)",
    "✅  Auto-generated PowerShell remediation scripts",
    "✅  Interactive Cytoscape.js graph visualization",
    "✅  REST API (FastAPI) with session management",
    "✅  Next.js dashboard with dark glassmorphism UI",
    "✅  Object Explorer, Attack Path Viewer, Findings Panel",
]

tb = add_text_box(slide, Inches(6.9), Inches(1.8), Inches(5.7), Inches(4.8))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "", font_size=1, color=BG_CARD)
for feat in features:
    add_paragraph(tf, feat, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT, space_before=Pt(3), space_after=Pt(1))

add_slide_number(slide, 10, TOTAL_SLIDES)

# =============================================================================
# SLIDE 11 — NEXT COURSE OF ACTION
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "Next Course of Action")
add_footer_line(slide)

# Left: Next steps
next_steps = [
    ("Q1", "GPO & Kerberos Deep Dive", "Extend enumeration to parse Group Policy Objects and Kerberos delegation settings for unconstrained/constrained delegation detection."),
    ("Q2", "Multi-Domain & Forest Support", "Scale the tool to handle multi-domain forests with cross-domain trust relationship analysis and inter-forest attack path discovery."),
    ("Q3", "Continuous Monitoring Mode", "Implement scheduled differential scans to detect configuration drift and alert on newly introduced misconfigurations in real-time."),
    ("Q4", "Reporting & Compliance", "Generate executive PDF reports, CSV exports, and compliance mappings (CIS Benchmarks, NIST 800-53) for audit and stakeholder communication."),
]

for i, (qtr, title, desc) in enumerate(next_steps):
    y = Inches(1.3) + i * Inches(1.35)
    # Quarter label
    qlabel = add_rect(slide, Inches(0.6), y + Inches(0.1), Inches(0.55), Inches(0.45), ACCENT_BLUE)
    qlabel.text_frame.paragraphs[0].text = qtr
    qlabel.text_frame.paragraphs[0].font.size = Pt(12)
    qlabel.text_frame.paragraphs[0].font.bold = True
    qlabel.text_frame.paragraphs[0].font.color.rgb = WHITE
    qlabel.text_frame.paragraphs[0].font.name = "Calibri"
    qlabel.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    card = add_rect(slide, Inches(1.3), y, Inches(11.5), Inches(1.15), BG_CARD, DARK_GRAY, Pt(0.75))
    tb = add_text_box(slide, Inches(1.5), y + Inches(0.08), Inches(11.1), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, title, font_size=15, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.LEFT)
    add_paragraph(tf, desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.JUSTIFY, space_before=Pt(4))

# Meeting schedule (right side or bottom)
tb = add_text_box(slide, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.5))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Guide Consultation:  Once a Week  (Every Monday/Tuesday)", font_size=13, bold=True, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 11, TOTAL_SLIDES)

# =============================================================================
# SLIDE 12 — GUIDE CONSULTATION + THANK YOU
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer_line(slide)

# Accent top bar
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_BLUE)

# Center card
card = add_rect(slide, Inches(2.5), Inches(1.5), Inches(8.333), Inches(4.5), BG_CARD, ACCENT_BLUE, Pt(1.5))

# Thank you
tb = add_text_box(slide, Inches(2.5), Inches(2.0), Inches(8.333), Inches(1.0))
set_text(tb.text_frame, "Thank You", font_size=44, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Divider
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.2), Inches(2.333), Pt(2))
div.fill.solid(); div.fill.fore_color.rgb = ACCENT_CYAN
div.line.fill.background(); div.shadow.inherit = False

# Questions
tb = add_text_box(slide, Inches(2.5), Inches(3.5), Inches(8.333), Inches(0.6))
set_text(tb.text_frame, "Questions & Discussion", font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Contact info
tb = add_text_box(slide, Inches(2.5), Inches(4.3), Inches(8.333), Inches(1.0))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Student Name  •  student.email@nfsu.ac.in", font_size=14, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "School of Cyber Security & Digital Forensics", font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(6))
add_paragraph(tf, "National Forensic Sciences University", font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(2))

# Bottom bar
add_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), ACCENT_BLUE)

add_slide_number(slide, 12, TOTAL_SLIDES)

# =============================================================================
# SAVE
# =============================================================================
prs.save(OUTPUT_FILE)
print(f"\n✅ Presentation saved to: {OUTPUT_FILE}")
print(f"   Total slides: {TOTAL_SLIDES}")
